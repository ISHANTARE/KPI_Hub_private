import collections 
import collections.abc
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pathlib import Path
import pandas as pd

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    title_shape.text = title
    subtitle_shape.text = subtitle
    return slide

def add_bullet_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = title
    tf = body_shape.text_frame
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            tf.text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            
    return slide

def main():
    prs = Presentation()
    
    # Slide 1: Title
    add_title_slide(prs, "KPI Hub Dashboard", "Engineering PMO - Executive Overview")
    
    # Slide 2: Executive Summary
    add_bullet_slide(prs, "Executive Summary", [
        "Comprehensive tracking of all active projects and components.",
        "Real-time visibility into portfolio health, schedule status, and critical risks.",
        "Integration with Codebeamer, Jira, and Azure DevOps for automated data aggregation.",
        "Predictive analytics for early detection of bottlenecks."
    ])
    
    # Slide 3: Key Performance Indicators (KPIs)
    add_bullet_slide(prs, "Key Performance Indicators (KPIs)", [
        "Portfolio Health Score: Continuous monitoring of overall portfolio wellness.",
        "On-Time Delivery Rate: Tracking schedule adherence across all domains.",
        "Quality Metrics: Defect density and test pass rates.",
        "Release Readiness: Combined assessment of health, quality, and outstanding risks."
    ])
    
    # Slide 4: Risk Management & Resource Utilization
    add_bullet_slide(prs, "Risk & Resource Management", [
        "Automated Risk Tracking: Highlighting critical risks with exposure scores.",
        "Resource Allocation: Visibility into overallocated and underutilized teams.",
        "Proactive Mitigation: Automated email notifications for resource bottlenecks.",
        "Traceability Insights: Ensuring testing coverage aligns with requirements."
    ])
    
    # Slide 5: AI-Powered Insights (Copilot)
    add_bullet_slide(prs, "AI-Powered Insights", [
        "Automated Weekly Summaries: AI-generated executive updates.",
        "Trend Analysis: Identifying underlying patterns in defects and schedule delays.",
        "Smart Recommendations: Actionable mitigation strategies suggested by the Copilot.",
        "Data-Driven Decisions: Empowering leadership with synthesized intelligence."
    ])
    
    # Slide 6: Next Steps
    add_bullet_slide(prs, "Next Steps & Action Items", [
        "Review critical risks flagged on the dashboard and assign mitigation owners.",
        "Rebalance overallocated resources based on the latest utilization report.",
        "Ensure all integrations (Jira, Codebeamer, etc.) are actively syncing.",
        "Enable weekly automated summary emails for the steering committee."
    ])
    
    out_path = Path("KPI_Hub_Dashboard_Overview.pptx")
    prs.save(out_path)
    print(f"Presentation saved successfully to {out_path.absolute()}")

if __name__ == "__main__":
    main()
