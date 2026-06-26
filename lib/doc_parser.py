import io
import pandas as pd

def extract_text_from_file(uploaded_file, file_type: str) -> str:
    """Extracts text from uploaded pptx, xlsx, or txt files."""
    text_content = ""
    try:
        if file_type == "txt":
            text_content = uploaded_file.read().decode("utf-8", errors="ignore")
            
        elif file_type == "eml":
            import email
            from email import policy
            msg = email.message_from_bytes(uploaded_file.read(), policy=policy.default)
            text_content = f"Subject: {msg['subject']}\nFrom: {msg['from']}\nTo: {msg['to']}\n\n"
            body = msg.get_body(preferencelist=('plain', 'html'))
            if body:
                text_content += body.get_content()
                
        elif file_type == "xlsx":
            # Just read all sheets using pandas and convert to string
            df_dict = pd.read_excel(uploaded_file, sheet_name=None)
            for sheet_name, df in df_dict.items():
                text_content += f"--- Sheet: {sheet_name} ---\n"
                text_content += df.to_string(index=False) + "\n\n"
                
        elif file_type == "pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(uploaded_file.read()))
            for i, slide in enumerate(prs.slides):
                text_content += f"--- Slide {i+1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"
                text_content += "\n"
    except Exception as e:
        return f"Error extracting text: {str(e)}"
        
    # Truncate if too long (keep within reasonable token limit)
    return text_content[:15000]
